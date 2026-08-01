"""
文档加载与章节感知分块。

对应 bRAG Notebook 1 的加载逻辑，改造点：
1. 用 book_id 关联 BookMetadata，而不是零散地传文件路径
2. PDF/EPUB 解析时提取章节标题，写入每个 chunk 的 chapter_title/chapter_index，
   而不是简单按页码切分——这是"单书精读"能精确定位到章节的前提

文档类型扩展（docx/pptx/html）：
沿用与 EPUB 相同的思路——直接读取文档自身的结构化标记（Word 标题样式 / PPT 每页 slide /
HTML 标题标签）作为章节边界，不做 OCR、不识别图片内容。这类解析器只适用于"文字层清晰、
可直接选中复制"的原生电子文档；文档中的扫描件、截图、图表贴图等图片内容目前不会被提取，
是已知能力边界，后续如需覆盖需引入 OCR/版面分析能力（如 MinerU）。

第 1 项输入处理优化（v1.6）：
1. 句子感知切分：RecursiveCharacterTextSplitter 的 separators 改为优先在中英文句末标点切割，
   避免句子中间断开导致的语义破碎。
2. 文本噪声清洗：_clean_text() 在送入切分器前对每个章节文本做统一清洗（连字符换行合并、
   多余空行折叠、页码行过滤、HTML 实体替换），减少无意义内容占据 chunk 空间。
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.core.config import get_settings
from app.models.schemas import BookFormat, BookMetadata, ChunkMetadata

# 常见章节标题模式：第X章、Chapter X、序、引言 等
_CHAPTER_PATTERNS = [
    re.compile(r"^第[0-9一二三四五六七八九十百千]+[章节回]\s*.*$"),
    re.compile(r"^Chapter\s+\d+.*$", re.IGNORECASE),
    re.compile(r"^\s*(序言|引言|前言|后记|附录)\s*.*$"),
]

# 句子感知切分分隔符列表（优先级从高到低）
# RecursiveCharacterTextSplitter 按列表顺序尝试分隔，优先在高层级边界切割，
# 字符数超限时才降级到下一个分隔符，最终兜底为单字符级别。
_SENTENCE_SEPARATORS = [
    "\n\n",       # 段落边界，最优先
    "。",          # 中文句末
    "！",
    "？",
    "…",          # 省略号（省略号后通常也是句子边界）
    ". ",          # 英文句末（带空格，避免将"3.14"等小数点误切）
    "! ",
    "? ",
    "\n",          # 换行
    "；", ";",     # 分号
    "，", ",",     # 逗号（最低优先，尽量避免在短语中间切）
    " ", "",       # 最后兜底
]

# 需要过滤的页码行模式（PDF 常见页眉页脚形式）
_PAGE_PATTERNS = [
    re.compile(r"\n第\s*\d+\s*页\n", re.IGNORECASE),
    re.compile(r"\nPage\s+\d+\n", re.IGNORECASE),
    re.compile(r"\n-\s*\d+\s*-\n"),  # "- 12 -" 形式
]

# HTML 实体映射（只处理最常见的几种，不引入 html.unescape 是为了避免把有意保留的实体也替换掉）
_HTML_ENTITIES = {
    "&amp;": "&",
    "&nbsp;": " ",
    "&lt;": "<",
    "&gt;": ">",
    "&quot;": '"',
    "&#39;": "'",
}


def _clean_text(text: str) -> str:
    """对章节原始文本做轻量噪声清洗，在送入切分器前统一调用。

    清洗规则（按顺序执行）：
    1. 英文连字符换行合并：PDF 排版常见，"recom-\\nmend" → "recommend"
    2. 多余空行折叠：3+ 连续空行压缩为 1 个空行，减少空白占用 chunk 空间
    3. 页码行过滤：移除常见"第 N 页"/"Page N"/"- N -"格式的页码行
    4. HTML 实体替换：处理 EPUB/HTML 文档遗留的 &amp; &nbsp; 等实体
    5. 首尾空白清理
    """
    # 1. 连字符换行合并
    text = re.sub(r"-\n", "", text)
    # 2. 多余空行折叠
    text = re.sub(r"\n{3,}", "\n\n", text)
    # 3. 页码行过滤
    for pattern in _PAGE_PATTERNS:
        text = pattern.sub("\n", text)
    # 4. HTML 实体替换
    for entity, char in _HTML_ENTITIES.items():
        text = text.replace(entity, char)
    # 5. 首尾空白
    return text.strip()


@dataclass
class ChapterBlock:
    index: int
    title: str | None
    page_start: int | None
    text: str


def _is_chapter_heading(line: str) -> bool:
    line = line.strip()
    if not line or len(line) > 40:
        return False
    return any(p.match(line) for p in _CHAPTER_PATTERNS)


def load_pdf_by_chapter(file_path: str) -> list[ChapterBlock]:
    """按章节切分 PDF。逐页扫描，命中章节标题时开启新 ChapterBlock。"""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    blocks: list[ChapterBlock] = []
    current = ChapterBlock(index=0, title=None, page_start=0, text="")

    for page_no, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        for line in page_text.splitlines():
            if _is_chapter_heading(line):
                if current.text.strip():
                    blocks.append(current)
                current = ChapterBlock(
                    index=len(blocks),
                    title=line.strip(),
                    page_start=page_no,
                    text="",
                )
            current.text += line + "\n"

    if current.text.strip():
        blocks.append(current)

    return blocks or [ChapterBlock(index=0, title=None, page_start=0, text=_pdf_fallback_text(reader))]


def _pdf_fallback_text(reader) -> str:
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def load_epub_by_chapter(file_path: str) -> list[ChapterBlock]:
    """按 EPUB 自身的 spine 顺序切分章节（每个 item 视为一章）。"""
    import ebooklib
    from bs4 import BeautifulSoup
    from ebooklib import epub

    book = epub.read_epub(file_path)
    blocks: list[ChapterBlock] = []

    for idx, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT)):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        title_tag = soup.find(["h1", "h2", "title"])
        title = title_tag.get_text(strip=True) if title_tag else None
        text = soup.get_text("\n").strip()
        if not text:
            continue
        blocks.append(ChapterBlock(index=len(blocks), title=title, page_start=None, text=text))

    return blocks


def load_txt_by_chapter(file_path: str) -> list[ChapterBlock]:
    """纯文本按章节标题行切分，没有标题则整本作为一个 block。"""
    content = Path(file_path).read_text(encoding="utf-8")
    lines = content.splitlines()
    blocks: list[ChapterBlock] = []
    current = ChapterBlock(index=0, title=None, page_start=None, text="")

    for line in lines:
        if _is_chapter_heading(line):
            if current.text.strip():
                blocks.append(current)
            current = ChapterBlock(index=len(blocks), title=line.strip(), page_start=None, text="")
        current.text += line + "\n"

    if current.text.strip():
        blocks.append(current)

    return blocks or [ChapterBlock(index=0, title=None, page_start=None, text=content)]


def load_docx_by_chapter(file_path: str) -> list[ChapterBlock]:
    """按 Word 大纲级别（outline level，对应"标题1/标题2..."样式）切分章节。

    用 outline_level 而不是 style.name 前缀匹配，因为 Word 内置标题样式名称是本地化的
    （中文"标题 1"、英文"Heading 1"、法文"Titre 1"等），outline_level 是底层 OOXML
    的 w:outlineLvl 属性，不受语言/自定义样式重命名影响，识别更稳定。

    只读取段落文字层，不解析文档中嵌入的图片；扫描件/截图等图片内容不会被提取。
    """
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    blocks: list[ChapterBlock] = []
    current = ChapterBlock(index=0, title=None, page_start=None, text="")

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        is_heading = paragraph.style.name != "Normal" and paragraph.paragraph_format.outline_level is not None
        if is_heading:
            if current.text.strip():
                blocks.append(current)
            current = ChapterBlock(index=len(blocks), title=text, page_start=None, text="")
        else:
            current.text += text + "\n"

    if current.text.strip():
        blocks.append(current)

    return blocks or [ChapterBlock(index=0, title=None, page_start=None, text=_docx_fallback_text(doc))]


def _docx_fallback_text(doc) -> str:
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_pptx_by_chapter(file_path: str) -> list[ChapterBlock]:
    """每一页 slide 作为一个 ChapterBlock，标题占位符文本作为章节标题。

    只读取 slide 中的文本框/占位符文字，不解析嵌入图片、图表内容；
    图片形式的图表/截图文字不会被提取。
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    blocks: list[ChapterBlock] = []

    for idx, slide in enumerate(prs.slides):
        title_shape = slide.shapes.title
        title = None
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text_frame.text.strip()
            if not shape_text:
                continue
            if shape is title_shape:
                title = shape_text
            else:
                texts.append(shape_text)

        text = "\n".join(texts)
        if not text and not title:
            continue
        blocks.append(ChapterBlock(index=len(blocks), title=title, page_start=idx, text=text))

    return blocks or [ChapterBlock(index=0, title=None, page_start=0, text="")]


def load_html_by_chapter(file_path: str) -> list[ChapterBlock]:
    """按 h1/h2 标签切分章节，静态解析 HTML 源码文字层（不执行 JS，不识别图片内容）。

    对每个标题标签向后收集同级兄弟节点文本，直到下一个标题标签为止；
    若文档不是"标题+同级正文"的扁平结构（如标题和正文分别嵌套在不同的 div 中），
    可能无法正确关联到对应正文，这是当前实现的已知局限。
    """
    from bs4 import BeautifulSoup

    content = Path(file_path).read_text(encoding="utf-8")
    soup = BeautifulSoup(content, "html.parser")

    headings = soup.find_all(["h1", "h2"])
    if not headings:
        text = soup.get_text("\n").strip()
        return [ChapterBlock(index=0, title=None, page_start=None, text=text)]

    blocks: list[ChapterBlock] = []
    for idx, heading in enumerate(headings):
        title = heading.get_text(strip=True)
        texts = []
        for sibling in heading.find_next_siblings():
            if sibling.name in ("h1", "h2"):
                break
            texts.append(sibling.get_text("\n", strip=True))
        text = "\n".join(t for t in texts if t)
        blocks.append(ChapterBlock(index=idx, title=title, page_start=None, text=text))

    return blocks


_LOADERS = {
    BookFormat.PDF: load_pdf_by_chapter,
    BookFormat.EPUB: load_epub_by_chapter,
    BookFormat.TXT: load_txt_by_chapter,
    BookFormat.DOCX: load_docx_by_chapter,
    BookFormat.PPTX: load_pptx_by_chapter,
    BookFormat.HTML: load_html_by_chapter,
}


def load_book_chapters(book: BookMetadata) -> list[ChapterBlock]:
    loader = _LOADERS.get(book.format)
    if loader is None:
        raise ValueError(f"不支持的书籍格式: {book.format}")
    return loader(book.source_path)


def _split_book_with_splitter(book: BookMetadata, splitter: RecursiveCharacterTextSplitter) -> list[Document]:
    """用给定的 splitter 把书按章节切分成 Document 列表，供不同粒度（单向量/父块）复用。

    每个章节文本在切分前先经过 _clean_text() 清洗，减少页码行、多余空行、HTML 实体等噪声。
    """
    chapters = load_book_chapters(book)
    documents: list[Document] = []
    chunk_seq = 0

    for chapter in chapters:
        cleaned = _clean_text(chapter.text)
        for chunk_text in splitter.split_text(cleaned):
            # 过滤空白 chunk 和极短 chunk（低于 20 字符无检索价值，会拉低 embedding 质量）
            if len(chunk_text.strip()) < 20:
                continue
            metadata = ChunkMetadata(
                book_id=book.book_id,
                book_title=book.title,
                chapter_index=chapter.index,
                chapter_title=chapter.title,
                page=chapter.page_start,
                chunk_index=chunk_seq,
            )
            documents.append(Document(page_content=chunk_text, metadata=metadata.to_payload()))
            chunk_seq += 1

    return documents


def split_book_into_documents(book: BookMetadata) -> list[Document]:
    """将一本书解析为带 ChunkMetadata 的 LangChain Document 列表。"""
    settings = get_settings()
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=_SENTENCE_SEPARATORS,
    )
    return _split_book_with_splitter(book, splitter)
